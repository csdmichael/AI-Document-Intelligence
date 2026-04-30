"""
Generate 100 PowerPoint presentations covering engineering filter design.

Each file contains 6 slides:
  1. Title — filter name, type, order, key spec
  2. Theory — transfer function, design equations, parameter table
  3. Frequency Response — matplotlib Bode plot (magnitude + phase) embedded as image
  4. Circuit/Block Diagram — matplotlib-drawn schematic embedded as image
  5. Design Parameters — 2-column table of spec → value
  6. Applications — bulleted real-world use cases

Output: data/pptx/filter_design_<family>_<variant>_<index>.pptx  (100 files total)
"""

import io
import os
import sys
import math

import matplotlib
matplotlib.use("Agg")  # non-interactive backend
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))
from config import cfg

# ---------------------------------------------------------------------------
# Output directory
# ---------------------------------------------------------------------------
OUTPUT_DIR = os.path.join(
    os.path.dirname(os.path.dirname(__file__)),
    cfg.app.pptx.output_dir,
)

# ---------------------------------------------------------------------------
# Filter definitions — 10 families × 10 variants = 100 presentations
# ---------------------------------------------------------------------------
FILTER_FAMILIES = [
    # (family_key, family_label, variants)
    ("butterworth_lp", "Butterworth Low-Pass", [
        {"order": 2, "fc": 1000, "ripple": None},
        {"order": 4, "fc": 2000, "ripple": None},
        {"order": 6, "fc": 5000, "ripple": None},
        {"order": 8, "fc": 10000, "ripple": None},
        {"order": 3, "fc": 500,  "ripple": None},
        {"order": 5, "fc": 3000, "ripple": None},
        {"order": 7, "fc": 8000, "ripple": None},
        {"order": 2, "fc": 15000, "ripple": None},
        {"order": 4, "fc": 100,  "ripple": None},
        {"order": 6, "fc": 20000, "ripple": None},
    ]),
    ("butterworth_hp", "Butterworth High-Pass", [
        {"order": 2, "fc": 1000, "ripple": None},
        {"order": 4, "fc": 500,  "ripple": None},
        {"order": 6, "fc": 2000, "ripple": None},
        {"order": 3, "fc": 5000, "ripple": None},
        {"order": 5, "fc": 200,  "ripple": None},
        {"order": 2, "fc": 8000, "ripple": None},
        {"order": 4, "fc": 100,  "ripple": None},
        {"order": 7, "fc": 3000, "ripple": None},
        {"order": 8, "fc": 12000, "ripple": None},
        {"order": 6, "fc": 750,  "ripple": None},
    ]),
    ("chebyshev1_lp", "Chebyshev Type-I Low-Pass", [
        {"order": 3, "fc": 1000, "ripple": 0.5},
        {"order": 5, "fc": 2000, "ripple": 1.0},
        {"order": 4, "fc": 500,  "ripple": 0.1},
        {"order": 6, "fc": 5000, "ripple": 3.0},
        {"order": 2, "fc": 10000, "ripple": 0.5},
        {"order": 7, "fc": 800,  "ripple": 1.0},
        {"order": 3, "fc": 3000, "ripple": 2.0},
        {"order": 5, "fc": 250,  "ripple": 0.5},
        {"order": 4, "fc": 15000, "ripple": 1.0},
        {"order": 8, "fc": 4000, "ripple": 0.5},
    ]),
    ("chebyshev2_lp", "Chebyshev Type-II Low-Pass", [
        {"order": 3, "fc": 1000, "ripple": 40},
        {"order": 5, "fc": 2000, "ripple": 60},
        {"order": 4, "fc": 500,  "ripple": 50},
        {"order": 6, "fc": 5000, "ripple": 40},
        {"order": 2, "fc": 10000, "ripple": 60},
        {"order": 7, "fc": 800,  "ripple": 40},
        {"order": 3, "fc": 3000, "ripple": 50},
        {"order": 5, "fc": 250,  "ripple": 60},
        {"order": 4, "fc": 15000, "ripple": 40},
        {"order": 8, "fc": 4000, "ripple": 50},
    ]),
    ("bessel_lp", "Bessel Low-Pass", [
        {"order": 2, "fc": 1000, "ripple": None},
        {"order": 4, "fc": 2000, "ripple": None},
        {"order": 3, "fc": 500,  "ripple": None},
        {"order": 5, "fc": 5000, "ripple": None},
        {"order": 6, "fc": 10000, "ripple": None},
        {"order": 2, "fc": 100,  "ripple": None},
        {"order": 4, "fc": 8000, "ripple": None},
        {"order": 7, "fc": 3000, "ripple": None},
        {"order": 3, "fc": 15000, "ripple": None},
        {"order": 5, "fc": 750,  "ripple": None},
    ]),
    ("elliptic_lp", "Elliptic Low-Pass", [
        {"order": 3, "fc": 1000, "ripple": 0.5},
        {"order": 5, "fc": 2000, "ripple": 1.0},
        {"order": 4, "fc": 500,  "ripple": 0.5},
        {"order": 6, "fc": 5000, "ripple": 1.0},
        {"order": 3, "fc": 10000, "ripple": 0.5},
        {"order": 5, "fc": 800,  "ripple": 1.0},
        {"order": 4, "fc": 3000, "ripple": 0.5},
        {"order": 6, "fc": 250,  "ripple": 1.0},
        {"order": 3, "fc": 15000, "ripple": 0.5},
        {"order": 5, "fc": 4000, "ripple": 1.0},
    ]),
    ("active_rc_lp", "Active RC Low-Pass", [
        {"order": 2, "fc": 1000, "ripple": None},
        {"order": 2, "fc": 5000, "ripple": None},
        {"order": 4, "fc": 500,  "ripple": None},
        {"order": 4, "fc": 2000, "ripple": None},
        {"order": 2, "fc": 10000, "ripple": None},
        {"order": 4, "fc": 100,  "ripple": None},
        {"order": 2, "fc": 20000, "ripple": None},
        {"order": 4, "fc": 800,  "ripple": None},
        {"order": 2, "fc": 3000, "ripple": None},
        {"order": 4, "fc": 15000, "ripple": None},
    ]),
    ("active_rc_bp", "Active RC Band-Pass", [
        {"order": 2, "fc": 1000, "ripple": None},
        {"order": 2, "fc": 5000, "ripple": None},
        {"order": 4, "fc": 500,  "ripple": None},
        {"order": 4, "fc": 2000, "ripple": None},
        {"order": 2, "fc": 10000, "ripple": None},
        {"order": 4, "fc": 100,  "ripple": None},
        {"order": 2, "fc": 20000, "ripple": None},
        {"order": 4, "fc": 800,  "ripple": None},
        {"order": 2, "fc": 3000, "ripple": None},
        {"order": 4, "fc": 15000, "ripple": None},
    ]),
    ("digital_iir", "Digital IIR Filter", [
        {"order": 4, "fc": 1000, "ripple": None},
        {"order": 6, "fc": 2000, "ripple": None},
        {"order": 4, "fc": 500,  "ripple": None},
        {"order": 8, "fc": 5000, "ripple": None},
        {"order": 4, "fc": 10000, "ripple": None},
        {"order": 6, "fc": 300,  "ripple": None},
        {"order": 8, "fc": 8000, "ripple": None},
        {"order": 4, "fc": 3000, "ripple": None},
        {"order": 6, "fc": 15000, "ripple": None},
        {"order": 8, "fc": 750,  "ripple": None},
    ]),
    ("digital_fir", "Digital FIR Filter", [
        {"order": 32,  "fc": 1000,  "ripple": None},
        {"order": 64,  "fc": 2000,  "ripple": None},
        {"order": 128, "fc": 500,   "ripple": None},
        {"order": 32,  "fc": 5000,  "ripple": None},
        {"order": 64,  "fc": 10000, "ripple": None},
        {"order": 128, "fc": 300,   "ripple": None},
        {"order": 32,  "fc": 8000,  "ripple": None},
        {"order": 64,  "fc": 3000,  "ripple": None},
        {"order": 128, "fc": 15000, "ripple": None},
        {"order": 32,  "fc": 750,   "ripple": None},
    ]),
]

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x1A, 0x23, 0x7E)
MID_BLUE  = RGBColor(0x15, 0x65, 0xC0)
LIGHT_BG  = RGBColor(0xE3, 0xF2, 0xFD)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT    = RGBColor(0xF9, 0xA8, 0x25)


# ---------------------------------------------------------------------------
# Helper: add a text box to a slide
# ---------------------------------------------------------------------------
def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=None, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return txBox


def _set_slide_background(slide, r, g, b):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def _add_title_slide(prs, family_label, variant, index_label):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide, 0x1A, 0x23, 0x7E)

    # Title
    _add_textbox(slide, 0.4, 0.8, 9.0, 1.2,
                 family_label,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    fc = variant["fc"]
    order = variant["order"]
    ripple = variant.get("ripple")
    ripple_str = f"  |  Passband Ripple: {ripple} dB" if ripple else ""
    subtitle = f"Order {order}  |  Cutoff Frequency: {fc:,} Hz{ripple_str}"
    _add_textbox(slide, 0.4, 2.2, 9.0, 0.6,
                 subtitle, font_size=18, color=LIGHT_BG, align=PP_ALIGN.CENTER)

    _add_textbox(slide, 0.4, 3.0, 9.0, 0.5,
                 index_label, font_size=13, color=ACCENT, align=PP_ALIGN.CENTER)

    # Decorative line
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        Inches(1.0), Inches(2.0), Inches(9.0), Inches(2.0),
    )
    line.line.color.rgb = ACCENT
    line.line.width = Pt(2)


# ---------------------------------------------------------------------------
# Slide 2 — Theory
# ---------------------------------------------------------------------------
def _add_theory_slide(prs, family_key, family_label, variant):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide, 0xFA, 0xFA, 0xFA)

    _add_textbox(slide, 0.3, 0.1, 9.3, 0.55,
                 "Filter Theory & Transfer Function",
                 font_size=20, bold=True, color=MID_BLUE)

    order = variant["order"]
    fc    = variant["fc"]
    wc    = 2 * math.pi * fc

    if "butterworth" in family_key:
        tf_text = (
            f"Transfer Function H(s) — Butterworth Order {order}\n\n"
            f"  |H(jω)|² = 1 / [1 + (ω/ωc)^{2*order}]\n\n"
            f"  ωc = 2π × {fc:,} rad/s  ≈  {wc:.1f} rad/s\n\n"
            f"  3 dB cutoff: {fc:,} Hz\n"
            f"  Maximally flat passband (no ripple)\n"
            f"  Roll-off: {20*order} dB/decade beyond ωc"
        )
    elif "chebyshev1" in family_key:
        ripple = variant.get("ripple", 1.0)
        tf_text = (
            f"Transfer Function H(s) — Chebyshev Type I Order {order}\n\n"
            f"  |H(jω)|² = 1 / [1 + ε²·Tₙ²(ω/ωc)]\n\n"
            f"  ε² = 10^(Rp/10) - 1,  Rp = {ripple} dB\n"
            f"  Tₙ = Chebyshev polynomial of order {order}\n"
            f"  ωc = 2π × {fc:,} rad/s\n"
            f"  Passband ripple: {ripple} dB\n"
            f"  Steeper roll-off than Butterworth for same order"
        )
    elif "chebyshev2" in family_key:
        ripple = variant.get("ripple", 40)
        tf_text = (
            f"Transfer Function H(s) — Chebyshev Type II Order {order}\n\n"
            f"  |H(jω)|² = 1 / [1 + 1/(ε²·Tₙ²(ωs/ω))]\n\n"
            f"  Stopband attenuation: {ripple} dB\n"
            f"  ωs = 2π × {fc:,} rad/s  (stopband edge)\n"
            f"  Flat passband, equiripple stopband\n"
            f"  Order: {order}"
        )
    elif "bessel" in family_key:
        tf_text = (
            f"Transfer Function H(s) — Bessel Order {order}\n\n"
            f"  H(s) = d₀ / Bₙ(s/ωc)\n\n"
            f"  Bₙ = Bessel polynomial of order {order}\n"
            f"  ωc = 2π × {fc:,} rad/s\n"
            f"  Maximally flat group delay\n"
            f"  Preserves waveform shape in passband\n"
            f"  Roll-off: slower than Butterworth"
        )
    elif "elliptic" in family_key:
        ripple = variant.get("ripple", 1.0)
        tf_text = (
            f"Transfer Function H(s) — Elliptic Order {order}\n\n"
            f"  |H(jω)|² = 1 / [1 + ε²·Rₙ²(ξ, ω/ωp)]\n\n"
            f"  ε² = 10^(Rp/10) - 1,  Rp = {ripple} dB\n"
            f"  Rₙ = Chebyshev rational function\n"
            f"  ωp = 2π × {fc:,} rad/s (passband edge)\n"
            f"  Equiripple in both passband and stopband\n"
            f"  Sharpest roll-off for given order and ripple"
        )
    elif "active_rc" in family_key and "bp" in family_key:
        tf_text = (
            f"Transfer Function H(s) — Active RC Band-Pass Order {order}\n\n"
            f"  H(s) = (s·BW) / [s² + s·(ω₀/Q) + ω₀²]\n\n"
            f"  ω₀ = 2π × {fc:,} rad/s  (center frequency)\n"
            f"  Q = ω₀/BW  (quality factor)\n"
            f"  Op-amp topology: multiple-feedback (MFB)\n"
            f"  Sections in cascade: {order // 2}"
        )
    elif "active_rc" in family_key:
        tf_text = (
            f"Transfer Function H(s) — Active RC Low-Pass Order {order}\n\n"
            f"  H(s) = ωc² / (s² + s·ωc/Q + ωc²)  [2nd-order section]\n\n"
            f"  ωc = 2π × {fc:,} rad/s\n"
            f"  Op-amp topology: Sallen-Key\n"
            f"  Number of 2nd-order stages: {order // 2}\n"
            f"  Unity-gain configuration (K=1)\n"
            f"  DC gain: 0 dB"
        )
    elif "digital_iir" in family_key:
        fs = fc * 10
        tf_text = (
            f"Transfer Function H(z) — Digital IIR Order {order}\n\n"
            f"  H(z) = B(z) / A(z)  =  Σbₖz⁻ᵏ / (1 + Σaₖz⁻ᵏ)\n\n"
            f"  Design: Bilinear transform from analog prototype\n"
            f"  Sampling rate: fs = {fs:,} Hz\n"
            f"  Normalized cutoff: ωd = 2π·{fc:,}/{fs:,} = {2*math.pi*fc/fs:.4f} rad\n"
            f"  Filter order: {order}  |  Coefficients: {order + 1}\n"
            f"  Recursive (feedback) structure"
        )
    else:  # FIR
        fs = fc * 10
        tf_text = (
            f"Transfer Function H(z) — Digital FIR Taps={order}\n\n"
            f"  H(z) = Σ h[n]·z⁻ⁿ,  n = 0…{order-1}\n\n"
            f"  Design: windowed-sinc (Hamming window)\n"
            f"  Sampling rate: fs = {fs:,} Hz\n"
            f"  Normalized cutoff: ωd = 2π·{fc:,}/{fs:,} = {2*math.pi*fc/fs:.4f} rad\n"
            f"  Non-recursive (FIR) — linear phase guaranteed\n"
            f"  Group delay: {order // 2} samples"
        )

    _add_textbox(slide, 0.4, 0.75, 9.0, 5.5,
                 tf_text, font_size=11, color=DARK_GRAY)


# ---------------------------------------------------------------------------
# Slide 3 — Frequency Response (Bode plot)
# ---------------------------------------------------------------------------
def _bode_magnitude(family_key, order, fc, frequencies):
    """Return magnitude response in dB for the given filter family."""
    wc = 2 * math.pi * fc
    mag = []
    for f in frequencies:
        w = 2 * math.pi * f
        ratio = w / wc

        try:
            if "butterworth" in family_key:
                h2 = 1.0 / (1.0 + ratio ** (2 * order))
            elif "chebyshev1" in family_key:
                ripple = 0.5
                eps2 = 10 ** (ripple / 10) - 1
                if ratio <= 1.0:
                    cheb = math.cos(order * math.acos(max(-1.0, min(1.0, ratio))))
                else:
                    cheb = math.cosh(order * math.acosh(ratio))
                h2 = 1.0 / (1.0 + eps2 * cheb ** 2)
            elif "chebyshev2" in family_key:
                if ratio == 0:
                    h2 = 1.0
                else:
                    inv = 1.0 / ratio
                    if inv <= 1.0:
                        cheb = math.cos(order * math.acos(max(-1.0, min(1.0, inv))))
                    else:
                        cheb = math.cosh(order * math.acosh(inv))
                    eps2 = 10 ** (40 / 10) - 1
                    h2 = 1.0 / (1.0 + 1.0 / (eps2 * cheb ** 2 + 1e-30))
            elif "bessel" in family_key:
                # Approximate Bessel as slightly softer Butterworth
                h2 = 1.0 / (1.0 + (ratio / 1.6) ** (2 * order))
            elif "elliptic" in family_key:
                # Approximate: steeper than Butterworth
                h2 = 1.0 / (1.0 + ratio ** (2 * order * 1.5))
            elif "bp" in family_key:
                # Band-pass centered at fc, BW = fc/2
                bw = fc / 2
                w0 = 2 * math.pi * fc
                if f == 0:
                    h2 = 0.0
                else:
                    q = w0 / (2 * math.pi * bw)
                    h2 = 1.0 / (1.0 + (q * (ratio - 1.0 / ratio)) ** (2 * max(order, 2)))
            elif "hp" in family_key:
                if ratio == 0:
                    h2 = 0.0
                else:
                    h2 = 1.0 / (1.0 + (1.0 / ratio) ** (2 * order))
            else:
                # Generic LP (covers FIR/IIR display)
                h2 = 1.0 / (1.0 + ratio ** (2 * min(order, 64)))
        except (OverflowError, ValueError, ZeroDivisionError):
            h2 = 0.0

        db = 10 * math.log10(max(h2, 1e-12))
        mag.append(db)
    return mag


def _add_frequency_response_slide(prs, family_key, family_label, variant):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide, 0xFA, 0xFA, 0xFA)

    _add_textbox(slide, 0.3, 0.05, 9.3, 0.5,
                 "Frequency Response — Bode Plot",
                 font_size=18, bold=True, color=MID_BLUE)

    fc    = variant["fc"]
    order = variant["order"]
    f_min = fc / 100
    f_max = fc * 100
    frequencies = np.logspace(math.log10(f_min), math.log10(f_max), 500)
    magnitude   = _bode_magnitude(family_key, order, fc, frequencies)

    # Phase (approximate)
    phase = []
    for f in frequencies:
        ratio = f / fc
        if "hp" in family_key:
            ph = 90 * order - order * math.degrees(math.atan(fc / f))
        elif "bp" in family_key:
            ph = math.degrees(math.atan(fc / (f + 1e-9))) - math.degrees(math.atan(f / (fc + 1e-9)))
        else:
            ph = -order * math.degrees(math.atan(f / fc))
        phase.append(ph)

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(8, 4.5), sharex=True)
    fig.patch.set_facecolor("#FAFAFA")

    ax1.semilogx(frequencies, magnitude, color="#1565C0", linewidth=2)
    ax1.set_ylabel("Magnitude (dB)", fontsize=9)
    ax1.set_title(f"{family_label} — Order {order} — fc={fc:,} Hz", fontsize=10)
    ax1.axvline(fc, color="#F9A825", linestyle="--", linewidth=1, label=f"fc={fc:,} Hz")
    ax1.axhline(-3, color="#C62828", linestyle=":", linewidth=1, label="-3 dB")
    ax1.legend(fontsize=7)
    ax1.grid(True, which="both", alpha=0.4)
    ax1.set_facecolor("#F5F9FF")

    ax2.semilogx(frequencies, phase, color="#2E7D32", linewidth=2)
    ax2.set_ylabel("Phase (°)", fontsize=9)
    ax2.set_xlabel("Frequency (Hz)", fontsize=9)
    ax2.axvline(fc, color="#F9A825", linestyle="--", linewidth=1)
    ax2.grid(True, which="both", alpha=0.4)
    ax2.set_facecolor("#F5FFF5")

    plt.tight_layout(pad=0.8)

    img_buf = io.BytesIO()
    fig.savefig(img_buf, format="png", dpi=120, bbox_inches="tight")
    plt.close(fig)
    img_buf.seek(0)

    slide.shapes.add_picture(img_buf, Inches(0.3), Inches(0.65), Inches(9.2), Inches(5.0))

    # Caption (alt-text / description)
    cap = (
        f"Bode plot of {family_label} order {order} filter. "
        f"Top panel: magnitude response in dB showing -3 dB cutoff at {fc:,} Hz "
        f"with a roll-off of approximately {20*order} dB/decade. "
        f"Bottom panel: phase response in degrees."
    )
    _add_textbox(slide, 0.3, 5.7, 9.2, 0.6, cap, font_size=8, color=DARK_GRAY)


# ---------------------------------------------------------------------------
# Slide 4 — Circuit / Block Diagram
# ---------------------------------------------------------------------------
def _add_circuit_slide(prs, family_key, family_label, variant):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide, 0xFA, 0xFA, 0xFA)

    _add_textbox(slide, 0.3, 0.05, 9.3, 0.5,
                 "Circuit / Block Diagram",
                 font_size=18, bold=True, color=MID_BLUE)

    order  = variant["order"]
    fc     = variant["fc"]
    fig, ax = plt.subplots(figsize=(8, 4.0))
    fig.patch.set_facecolor("#FAFAFA")
    ax.set_facecolor("#FAFAFA")
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 5)
    ax.axis("off")

    if "active_rc" in family_key or ("butterworth" in family_key and order <= 4):
        # Draw Sallen-Key 2nd-order LP section
        ax.set_title(f"Sallen-Key Low-Pass Stage — 2nd Order Section\n{family_label}", fontsize=10)
        # Op-amp symbol (triangle)
        amp_x, amp_y = 6.5, 2.5
        amp = mpatches.FancyArrow(amp_x, amp_y - 0.5, 0.0, 1.0, width=0.0, head_width=0.0,
                                   length_includes_head=True, color="none")
        triangle = plt.Polygon([[amp_x, amp_y - 0.6], [amp_x, amp_y + 0.6], [amp_x + 1.0, amp_y]],
                                closed=True, fill=False, edgecolor="#1565C0", linewidth=2)
        ax.add_patch(triangle)
        # + and - inputs
        ax.text(amp_x + 0.1, amp_y + 0.3, "+", fontsize=10, color="#1565C0", ha="center")
        ax.text(amp_x + 0.1, amp_y - 0.3, "−", fontsize=10, color="#1565C0", ha="center")
        # Input wire
        ax.annotate("", xy=(amp_x, amp_y + 0.3), xytext=(0.5, amp_y + 0.3),
                    arrowprops=dict(arrowstyle="->", color="black", lw=1.5))
        # R1 box
        r1 = mpatches.FancyBboxPatch((1.5, amp_y + 0.1), 1.2, 0.4, boxstyle="round,pad=0.05",
                                      linewidth=1.5, edgecolor="#333", facecolor="#E3F2FD")
        ax.add_patch(r1)
        ax.text(2.1, amp_y + 0.3, f"R₁", fontsize=9, ha="center", va="center")
        # R2 box
        r2 = mpatches.FancyBboxPatch((3.3, amp_y + 0.1), 1.2, 0.4, boxstyle="round,pad=0.05",
                                      linewidth=1.5, edgecolor="#333", facecolor="#E3F2FD")
        ax.add_patch(r2)
        ax.text(3.9, amp_y + 0.3, f"R₂", fontsize=9, ha="center", va="center")
        # C1 (shunt)
        ax.plot([2.7, 2.7], [amp_y + 0.1, amp_y - 0.8], color="black", lw=1.5)
        c1 = mpatches.FancyBboxPatch((2.35, amp_y - 1.2), 0.7, 0.4, boxstyle="round,pad=0.05",
                                      linewidth=1.5, edgecolor="#333", facecolor="#FFF9C4")
        ax.add_patch(c1)
        ax.text(2.7, amp_y - 1.0, "C₁", fontsize=9, ha="center", va="center")
        ax.plot([2.7, 2.7], [amp_y - 1.2, amp_y - 1.6], color="black", lw=1.5)
        ax.plot([2.0, 3.4], [amp_y - 1.6, amp_y - 1.6], color="black", lw=1.5)
        ax.text(2.7, amp_y - 1.9, "GND", fontsize=8, ha="center", color="#666")
        # C2 (feedback)
        ax.plot([5.5, 5.5, 7.5, 7.5], [amp_y + 0.3, amp_y + 1.4, amp_y + 1.4, amp_y],
                color="black", lw=1.5)
        c2 = mpatches.FancyBboxPatch((5.2, amp_y + 1.2), 0.7, 0.4, boxstyle="round,pad=0.05",
                                      linewidth=1.5, edgecolor="#333", facecolor="#FFF9C4")
        ax.add_patch(c2)
        ax.text(5.55, amp_y + 1.4, "C₂", fontsize=9, ha="center", va="center")
        # Output
        ax.annotate("", xy=(9.5, amp_y), xytext=(7.5, amp_y),
                    arrowprops=dict(arrowstyle="->", color="black", lw=1.5))
        ax.text(9.6, amp_y, "Vout", fontsize=9, va="center")
        ax.text(0.3, amp_y + 0.3, "Vin", fontsize=9, va="center")
        # Component labels
        wc = 2 * math.pi * fc
        ax.text(5.0, 0.3,
                f"ωc = 2π × {fc:,} Hz    R = {1/(wc*1e-9):.0f} Ω (ideal)    C = 1 nF (ideal)",
                fontsize=8, ha="center", color="#333")

    elif "fir" in family_key or "iir" in family_key:
        # Digital filter block diagram
        taps = min(order, 8)
        ax.set_title(f"Direct Form II Transposed — {'FIR' if 'fir' in family_key else 'IIR'} N={order}", fontsize=10)
        box_w, box_h = 0.7, 0.55
        gap = (9.0 - taps * box_w) / (taps + 1)
        y_reg = 3.0
        y_mult = 1.5
        x_positions = [gap + i * (box_w + gap) + 0.5 for i in range(taps)]
        # Input arrow
        ax.annotate("", xy=(x_positions[0] + box_w / 2, y_reg + box_h / 2),
                    xytext=(x_positions[0] + box_w / 2, y_reg + box_h + 0.5),
                    arrowprops=dict(arrowstyle="->", color="black", lw=1.5))
        ax.text(x_positions[0] + box_w / 2, y_reg + box_h + 0.65, "x[n]", ha="center", fontsize=9)
        for i, xp in enumerate(x_positions):
            # Delay register
            reg = mpatches.FancyBboxPatch((xp, y_reg), box_w, box_h, boxstyle="round,pad=0.05",
                                           linewidth=1.5, edgecolor="#1565C0", facecolor="#E3F2FD")
            ax.add_patch(reg)
            ax.text(xp + box_w / 2, y_reg + box_h / 2, f"z⁻¹", ha="center", va="center", fontsize=9)
            # Multiplier
            coeff_circle = plt.Circle((xp + box_w / 2, y_mult), 0.28,
                                       fill=True, facecolor="#FFF9C4", edgecolor="#333", linewidth=1.5)
            ax.add_patch(coeff_circle)
            ax.text(xp + box_w / 2, y_mult, f"h[{i}]", ha="center", va="center", fontsize=7.5)
            # Vertical line from register to multiplier
            ax.plot([xp + box_w / 2, xp + box_w / 2], [y_reg, y_mult + 0.28], color="black", lw=1.2)
            # Vertical line from multiplier to summing bus
            ax.plot([xp + box_w / 2, xp + box_w / 2], [y_mult - 0.28, 0.7], color="black", lw=1.2)
            # Connect delay registers in series
            if i < taps - 1:
                ax.annotate("", xy=(x_positions[i + 1], y_reg + box_h / 2),
                            xytext=(xp + box_w, y_reg + box_h / 2),
                            arrowprops=dict(arrowstyle="->", color="black", lw=1.5))
        # Summing bus
        ax.plot([x_positions[0] + box_w / 2, x_positions[-1] + box_w / 2], [0.7, 0.7],
                color="#C62828", lw=2.0)
        # Output
        mid_x = (x_positions[0] + x_positions[-1]) / 2 + box_w / 2
        ax.annotate("", xy=(mid_x, 0.1), xytext=(mid_x, 0.7),
                    arrowprops=dict(arrowstyle="->", color="#C62828", lw=1.8))
        ax.text(mid_x, 0.0, "y[n]", ha="center", fontsize=9, color="#C62828")
        ax.text(5.0, 4.7,
                f"fs = {fc*10:,} Hz  |  fc = {fc:,} Hz  |  Taps = {order}",
                fontsize=9, ha="center", color="#333")

    else:
        # Generic block diagram for other types
        ax.set_title(f"{family_label} — Signal Flow", fontsize=11)
        stages = max(1, order // 2)
        stage_w = 1.4
        gap_x = (9.0 - stages * stage_w) / (stages + 1)
        y_center = 2.5
        prev_x = 0.3
        ax.annotate("", xy=(prev_x + 0.5, y_center), xytext=(0.3, y_center),
                    arrowprops=dict(arrowstyle="->", color="black", lw=1.5))
        ax.text(0.1, y_center, "In", fontsize=10, va="center")
        for s in range(stages):
            bx = gap_x + s * (stage_w + gap_x) + 0.5
            stage_box = mpatches.FancyBboxPatch((bx, y_center - 0.5), stage_w, 1.0,
                                                 boxstyle="round,pad=0.05",
                                                 linewidth=2, edgecolor="#1565C0", facecolor="#E3F2FD")
            ax.add_patch(stage_box)
            ax.text(bx + stage_w / 2, y_center + 0.15, f"Stage {s+1}", ha="center", va="center",
                    fontsize=9, fontweight="bold", color="#1565C0")
            ax.text(bx + stage_w / 2, y_center - 0.2, f"H{s+1}(s)", ha="center", va="center",
                    fontsize=8.5, color="#333")
            if s < stages - 1:
                next_bx = gap_x + (s + 1) * (stage_w + gap_x) + 0.5
                ax.annotate("", xy=(next_bx, y_center), xytext=(bx + stage_w, y_center),
                            arrowprops=dict(arrowstyle="->", color="black", lw=1.5))
        last_bx = gap_x + (stages - 1) * (stage_w + gap_x) + 0.5 + stage_w
        ax.annotate("", xy=(9.7, y_center), xytext=(last_bx, y_center),
                    arrowprops=dict(arrowstyle="->", color="black", lw=1.5))
        ax.text(9.8, y_center, "Out", fontsize=10, va="center")
        ax.text(5.0, 0.3, f"fc = {fc:,} Hz  |  Order = {order}  |  {stages} cascaded section(s)",
                fontsize=9, ha="center", color="#333")

    plt.tight_layout(pad=0.5)
    img_buf = io.BytesIO()
    fig.savefig(img_buf, format="png", dpi=120, bbox_inches="tight")
    plt.close(fig)
    img_buf.seek(0)

    slide.shapes.add_picture(img_buf, Inches(0.3), Inches(0.65), Inches(9.2), Inches(4.8))

    # Caption
    if "active_rc" in family_key or "butterworth" in family_key:
        cap = (
            f"Sallen-Key topology for {family_label} order {order}. "
            f"The circuit uses two resistors R₁, R₂ and two capacitors C₁, C₂ "
            f"around an op-amp in unity-gain configuration to realise a 2nd-order "
            f"low-pass section with cutoff frequency {fc:,} Hz."
        )
    elif "fir" in family_key or "iir" in family_key:
        cap = (
            f"Direct Form II transposed signal-flow graph for a "
            f"{'FIR' if 'fir' in family_key else 'IIR'} filter of order {order}. "
            f"Each z⁻¹ block represents a one-sample delay register. "
            f"Coefficients h[n] are multiplied and summed to form the output y[n]."
        )
    else:
        stages = max(1, order // 2)
        cap = (
            f"Cascade block diagram for {family_label} order {order}. "
            f"The design is decomposed into {stages} second-order section(s), "
            f"each realising a pair of complex conjugate poles near cutoff {fc:,} Hz."
        )
    _add_textbox(slide, 0.3, 5.55, 9.2, 0.7, cap, font_size=8, color=DARK_GRAY)


# ---------------------------------------------------------------------------
# Slide 5 — Design Parameters Table
# ---------------------------------------------------------------------------
def _add_parameters_slide(prs, family_key, family_label, variant):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide, 0xFA, 0xFA, 0xFA)

    _add_textbox(slide, 0.3, 0.05, 9.3, 0.5,
                 "Design Parameters",
                 font_size=18, bold=True, color=MID_BLUE)

    order = variant["order"]
    fc    = variant["fc"]
    wc    = 2 * math.pi * fc
    ripple = variant.get("ripple")

    rows = [
        ("Parameter", "Value"),
        ("Filter Family", family_label),
        ("Filter Type", "Low-Pass" if "lp" in family_key else ("High-Pass" if "hp" in family_key else "Band-Pass")),
        ("Filter Order (N)", str(order)),
        ("Cutoff Frequency (fc)", f"{fc:,} Hz"),
        ("Angular Cutoff (ωc)", f"{wc:.2f} rad/s"),
        ("Passband Gain (0 Hz)", "0 dB"),
        ("-3 dB Frequency", f"{fc:,} Hz"),
        ("Roll-off Rate", f"{20 * order} dB/decade"),
    ]
    if ripple is not None:
        if "chebyshev2" in family_key:
            rows.append(("Minimum Stopband Attenuation", f"{ripple} dB"))
        else:
            rows.append(("Passband Ripple", f"{ripple} dB"))
    if "active_rc" in family_key:
        wc_val = 2 * math.pi * fc
        rows.append(("Op-Amp Topology", "Sallen-Key"))
        rows.append(("Ideal Capacitor (C)", "1 nF"))
        rows.append(("Ideal Resistor (R)", f"{1/(wc_val * 1e-9):.0f} Ω"))
    if "digital" in family_key:
        fs = fc * 10
        rows.append(("Sampling Rate (fs)", f"{fs:,} Hz"))
        rows.append(("Nyquist Frequency", f"{fs // 2:,} Hz"))
        rows.append(("Normalised Cutoff (ωd)", f"{2 * math.pi * fc / fs:.4f} rad/sample"))
    if "fir" in family_key:
        rows.append(("Window Function", "Hamming"))
        rows.append(("Linear Phase", "Yes"))
        rows.append(("Group Delay", f"{order // 2} samples"))

    tbl = slide.shapes.add_table(len(rows), 2, Inches(0.5), Inches(0.65),
                                  Inches(9.0), Inches(min(5.5, 0.4 * len(rows)))).table
    for row_i, (param, value) in enumerate(rows):
        cell0 = tbl.cell(row_i, 0)
        cell1 = tbl.cell(row_i, 1)
        cell0.text = param
        cell1.text = value

        # Style header row
        is_header = row_i == 0
        for cell in (cell0, cell1):
            para = cell.text_frame.paragraphs[0]
            run  = para.runs[0] if para.runs else para.add_run()
            run.font.size = Pt(11 if is_header else 10)
            run.font.bold = is_header
            if is_header:
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = MID_BLUE
            elif row_i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)


# ---------------------------------------------------------------------------
# Slide 6 — Applications
# ---------------------------------------------------------------------------
def _add_applications_slide(prs, family_key, family_label, variant):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_background(slide, 0xFA, 0xFA, 0xFA)

    _add_textbox(slide, 0.3, 0.05, 9.3, 0.5,
                 "Real-World Applications",
                 font_size=18, bold=True, color=MID_BLUE)

    fc = variant["fc"]

    app_map = {
        "butterworth_lp": [
            "• Anti-aliasing filter before ADC sampling in data acquisition systems",
            "• Audio signal smoothing in professional mixing consoles",
            "• DC power supply ripple rejection in instrumentation amplifiers",
            "• Vibration sensor conditioning in industrial IoT edge devices",
            "• Baseband filtering in software-defined radio (SDR) receivers",
        ],
        "butterworth_hp": [
            "• DC-blocking in microphone pre-amplifier circuits",
            "• Removal of low-frequency mechanical vibration interference",
            "• Coupling filter in AC-coupled differential amplifier inputs",
            "• High-pass shelving EQ in professional audio mastering",
            "• Seismic noise rejection in structural health monitoring",
        ],
        "chebyshev1_lp": [
            "• Steeper anti-aliasing where slight passband ripple is acceptable",
            "• Channel selection in narrow-band RF receiver front-ends",
            "• Harmonic rejection in precision measurement equipment",
            "• Video signal bandwidth limiting for broadcast standards",
            "• Medical ECG signal processing — power-line noise rejection",
        ],
        "chebyshev2_lp": [
            "• Stopband-critical filtering in telephony channel banks",
            "• Interference suppression in GPS front-end receiver chains",
            "• Image rejection in superheterodyne radio architectures",
            "• Flat passband required: medical imaging signal chains",
            "• High-fidelity audio crossover networks (subwoofer isolation)",
        ],
        "bessel_lp": [
            "• Pulse shaping in digital communications (minimum ISI)",
            "• Group-delay equalisation in multi-channel audio systems",
            "• Waveform preservation in oscilloscope input paths",
            "• Radar pulse processing requiring constant group delay",
            "• Control system compensation for minimum phase distortion",
        ],
        "elliptic_lp": [
            "• Narrowest transition band: satellite communication transponders",
            "• Multi-tone DTMF decoder front-end filtering",
            "• Anti-aliasing in high-speed ADCs where order must be minimised",
            "• Channel guard filtering in dense OFDM subcarrier spacing",
            "• Spectral containment in military frequency-hopping radios",
        ],
        "active_rc_lp": [
            f"• Instrumentation amplifier output smoothing at {fc:,} Hz",
            "• Sensor conditioning circuit for MEMS accelerometers",
            "• Loop filter in phase-locked loop (PLL) synthesisers",
            "• Anti-aliasing for low-speed (audio-band) delta-sigma ADCs",
            "• Low-pass reconstruction filter after DAC in control systems",
        ],
        "active_rc_bp": [
            f"• Bandpass channel filter centred at {fc:,} Hz in DTMF decoder",
            "• Resonant peak equaliser in parametric EQ systems",
            "• Lock-in amplifier reference channel selection",
            "• Carrier frequency selection in power-line communication",
            "• Biomedical EEG alpha/beta/gamma band extraction",
        ],
        "digital_iir": [
            "• Real-time audio effects processing on DSP microcontrollers",
            "• Recursive Butterworth implementation in FPGA signal chains",
            f"• {fc:,} Hz lowpass post-processing in embedded sensor fusion",
            "• Speech enhancement preprocessing in voice recognition engines",
            "• Power spectral density estimation in industrial vibration analysis",
        ],
        "digital_fir": [
            "• Linear-phase FIR for matched filtering in digital communications",
            "• Pulse-shaping (root-raised cosine) in QAM modem transmitters",
            f"• {fc:,} Hz decimation filter before sample-rate reduction",
            "• Noise-reduction preprocessing in hearing aid DSP chips",
            "• Anti-imaging reconstruction filter after DAC upsampling",
        ],
    }

    apps = app_map.get(family_key, [
        "• Signal conditioning in data acquisition systems",
        "• Noise rejection in analog front-end circuits",
        "• Interference suppression in communication receivers",
        "• Pre-processing for digital signal analysis",
        "• Bandwidth limitation for power-efficient transmission",
    ])

    _add_textbox(slide, 0.5, 0.65, 9.0, 4.8,
                 "\n".join(apps), font_size=14, color=DARK_GRAY)

    # Footer note
    _add_textbox(slide, 0.5, 5.6, 9.0, 0.55,
                 f"Filter: {family_label} | fc = {fc:,} Hz | Order = {variant['order']}",
                 font_size=9, color=MID_BLUE)


# ---------------------------------------------------------------------------
# Main: generate one PPTX per variant
# ---------------------------------------------------------------------------
def create_filter_pptx(family_key: str, family_label: str, variant: dict,
                        variant_index: int, global_index: int):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    index_label = f"File {global_index:03d} / 100  |  Variant {variant_index + 1}"

    _add_title_slide(prs, family_label, variant, index_label)
    _add_theory_slide(prs, family_key, family_label, variant)
    _add_frequency_response_slide(prs, family_key, family_label, variant)
    _add_circuit_slide(prs, family_key, family_label, variant)
    _add_parameters_slide(prs, family_key, family_label, variant)
    _add_applications_slide(prs, family_key, family_label, variant)

    filename = f"filter_design_{family_key}_{variant_index + 1:02d}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)
    return filename


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    total = 0
    for family_key, family_label, variants in FILTER_FAMILIES:
        for vi, variant in enumerate(variants):
            total += 1
            filename = create_filter_pptx(family_key, family_label, variant, vi, total)
            if total % 10 == 0:
                print(f"  Generated {total}/100 — {filename}")

    print(f"\nDone! Generated {total} PPTX filter-design files in {OUTPUT_DIR}")


if __name__ == "__main__":
    main()
